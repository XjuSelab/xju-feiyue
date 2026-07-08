"""生成《飞跃·笔记系统与社区互动模块》小学期答辩 PPT (.pptx)。

- 18 页，16:9；配 docs/coursework/figures/*.png
- 口径：最新完成度（治理后端已实现+212测试；前端已开发·本地验收中）
- 只讲产品内的 AI 摘要/润色能力，不涉及开发过程
"""
from __future__ import annotations

import struct
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "docs" / "coursework" / "figures"
OUT = ROOT / "docs" / "赵文彪-飞跃-小学期答辩.pptx"

# ---- 主题 ----
INK = RGBColor(0x1B, 0x1F, 0x24)
PRIMARY = RGBColor(0x1F, 0x4E, 0x79)      # 深蓝
ACCENT = RGBColor(0x2E, 0x86, 0xC1)       # 亮蓝
GRAY = RGBColor(0x5B, 0x63, 0x6B)
LIGHT = RGBColor(0xEE, 0xF2, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "微软雅黑"

SW, SH = Inches(13.333), Inches(7.5)


def png_size(p: Path) -> tuple[int, int]:
    with open(p, "rb") as f:
        head = f.read(24)
    return struct.unpack(">II", head[16:24])


def _set_font(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tb, tf


def _rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _para(tf, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT, space_after=8,
          first=False, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.level = level
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    _set_font(r)
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_header(slide, title, page):
    _rect(slide, 0, 0, SW, Inches(0.16), PRIMARY)
    _rect(slide, Inches(0.55), Inches(0.62), Inches(0.09), Inches(0.52), ACCENT)
    _, tf = _box(slide, Inches(0.78), Inches(0.52), Inches(11.6), Inches(0.8))
    _para(tf, title, 26, PRIMARY, bold=True, first=True, space_after=0)
    # footer
    _, ff = _box(slide, Inches(0.55), Inches(7.02), Inches(9), Inches(0.35))
    _para(ff, "飞跃平台 · 笔记系统与社区互动模块", 9, GRAY, first=True, space_after=0)
    _, pf = _box(slide, Inches(11.8), Inches(7.02), Inches(1.0), Inches(0.35))
    _para(pf, f"{page:02d}", 10, GRAY, align=PP_ALIGN.RIGHT, first=True, space_after=0)


def bullets(slide, items, x, y, w, h, size=16, gap=10):
    _, tf = _box(slide, x, y, w, h)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        p = _para(tf, ("▸ " if lvl == 0 else "· ") + txt, size - lvl * 2,
                  INK if lvl == 0 else GRAY, bold=(lvl == 0), first=(i == 0),
                  space_after=gap, level=lvl)


def figure(slide, img, x, y, maxw, maxh, caption=None):
    w, h = png_size(img)
    r = min(maxw / w, maxh / h)
    fw, fh = int(w * r), int(h * r)
    px = x + (maxw - fw) // 2
    py = y + (maxh - fh) // 2
    slide.shapes.add_picture(str(img), Emu(int(px)), Emu(int(py)), Emu(fw), Emu(fh))
    if caption:
        _, cf = _box(slide, x, Emu(int(y + maxh + Emu(0.02))), Emu(int(maxw)), Inches(0.3))
        _para(cf, caption, 10, GRAY, align=PP_ALIGN.CENTER, first=True, space_after=0)


# ==========================================================================
def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    # ---- 1 封面 ----
    s = blank(prs)
    _rect(s, 0, 0, SW, SH, PRIMARY)
    _rect(s, 0, Inches(4.75), SW, Inches(0.06), ACCENT)
    _, tf = _box(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.2))
    _para(tf, "飞跃平台", 30, RGBColor(0xBF, 0xD7, 0xEA), bold=True, first=True, space_after=6)
    _para(tf, "笔记系统与社区互动模块", 42, WHITE, bold=True, space_after=0)
    _, sf = _box(s, Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.6))
    _para(sf, "小学期课程设计 · 答辩汇报", 20, RGBColor(0xD9, 0xE6, 0xF2), first=True, space_after=14)
    _para(sf, "汇报人：赵文彪        2026 年 7 月", 16, RGBColor(0xBF, 0xD7, 0xEA), space_after=0)

    # ---- 2 目录 ----
    s = blank(prs)
    add_header(s, "目录", 2)
    cat = [
        ("01", "项目背景与需求范围"),
        ("02", "系统分析与总体设计"),
        ("03", "核心功能设计与实现"),
        ("04", "治理子系统与关键技术"),
        ("05", "测试与验证"),
        ("06", "总结与展望"),
    ]
    _, tf = _box(s, Inches(2.2), Inches(1.9), Inches(9), Inches(5))
    for i, (n, t) in enumerate(cat):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(16)
        r1 = p.add_run(); r1.text = n + "  "
        r1.font.size = Pt(24); r1.font.bold = True; r1.font.color.rgb = ACCENT; _set_font(r1)
        r2 = p.add_run(); r2.text = t
        r2.font.size = Pt(22); r2.font.color.rgb = INK; _set_font(r2)

    # ---- 3 背景与定位 ----
    s = blank(prs)
    add_header(s, "一、项目背景与模块定位", 3)
    bullets(s, [
        "飞跃平台面向高校学生的升学 / 竞赛 / 课程经验共享",
        "痛点：经验散落在群聊与口头交流 —— 难检索、难沉淀、难维护",
        "本模块职责：",
        ("内容生产（写作）· 内容展示（信息流 / 详情页）", 1),
        ("社区互动（评论 / 点赞点踩 / 收藏 / 表态）", 1),
        ("成长激励（每日签到 / 经验 / 等级）· 内容治理（举报 / 拉黑 / AI 审核）", 1),
        "定位：平台内容核心层，连接写作页、信息流、详情页、个人中心与后端数据模型",
    ], Inches(0.9), Inches(1.7), Inches(11.6), Inches(5))

    # ---- 4 需求范围·增量模型 ----
    s = blank(prs)
    add_header(s, "二、需求范围 · 增量模型", 4)
    bullets(s, [
        "增量一 · 内容：草稿自动保存、发布笔记、信息流浏览、详情阅读、AI 辅助写作与摘要",
        "增量二 · 互动与成长：点赞 / 点踩、收藏、评论与二级回复、评论表态、合集、每日签到与等级",
        "增量三 · 治理：举报、拉黑、AI 后台预审、人工裁决、操作频次限流",
        ("本阶段已落地：增量一 / 二核心链路 + 增量三治理后端（含 AI 预审）", 1),
        ("如实边界：治理与部分互动的前端入口已开发，正在本地验收接入", 1),
    ], Inches(0.9), Inches(1.7), Inches(11.6), Inches(5))

    # ---- 5 需求分析·用例图 ----
    s = blank(prs)
    add_header(s, "三、需求分析 · 用例图", 5)
    bullets(s, [
        "4 类执行者：匿名访客 / 注册用户 / 笔记作者 / 外部 AI 服务 · 共 18 个用例",
    ], Inches(0.9), Inches(1.55), Inches(11.6), Inches(0.6), size=15)
    figure(s, FIG / "ch4-fig5-usecase.png", Inches(1.2), Inches(2.15), Inches(11), Inches(4.55),
           caption="图 4-5  模块用例图")

    # ---- 6 数据流图 ----
    s = blank(prs)
    add_header(s, "三、需求分析 · 顶层数据流图", 6)
    bullets(s, [
        "标准 DFD 画法：矩形 = 外部实体，圆形 = 加工，开口矩形 = 数据存储",
    ], Inches(0.9), Inches(1.55), Inches(11.6), Inches(0.6), size=15)
    figure(s, FIG / "ch4-fig2-dfd-top.png", Inches(0.9), Inches(2.15), Inches(11.5), Inches(4.55),
           caption="图 4-2  顶层数据流图")

    # ---- 7 系统架构 ----
    s = blank(prs)
    add_header(s, "四、系统总体设计 · 分层架构", 7)
    bullets(s, [
        "前端 React 18 SPA · 后端 FastAPI 单体（routes → services → db）",
        "数据层 SQLite / SQLAlchemy ORM（模型为单一结构来源）",
        "AI 经 OpenAI 兼容协议接 DeepSeek，仅用于润色与摘要，测试 dry-run",
    ], Inches(0.9), Inches(1.6), Inches(4.7), Inches(5), size=15)
    figure(s, FIG / "ch5-fig1-module-structure.png", Inches(5.7), Inches(1.7), Inches(7.1),
           Inches(4.9), caption="图 5-1  模块结构图")

    # ---- 8 领域数据模型 ----
    s = blank(prs)
    add_header(s, "四、总体设计 · 领域数据模型（UML 类图）", 8)
    bullets(s, [
        "18 个实体：User / Note / Comment / Collection / Report / Block …",
        "三格类框：类名 · 属性 · 方法（已按代码实现替换 ER 图）",
    ], Inches(0.9), Inches(1.55), Inches(11.6), Inches(0.85), size=14)
    figure(s, FIG / "ch5-fig2-class-domain.png", Inches(0.9), Inches(2.35), Inches(11.5),
           Inches(4.4), caption="图 5-2  领域模型类图")

    # ---- 9 核心链路一 ----
    s = blank(prs)
    add_header(s, "五、核心链路一 · 草稿发布与 AI 摘要降级", 9)
    bullets(s, [
        "自动保存草稿 → 发布前校验标题 / 正文 / 分类",
        "摘要为空 → 调 AI 生成一句话简介",
        "AI 超时 / 失败 → 首段截断兜底",
        "单事务创建 Note 并删除 Draft",
        ("设计原则：AI 永不阻断主流程", 1),
    ], Inches(0.9), Inches(1.6), Inches(4.7), Inches(5), size=15)
    figure(s, FIG / "ch6-fig8-seq-publish.png", Inches(5.7), Inches(1.7), Inches(7.1),
           Inches(4.9), caption="图 6-8  发布顺序图")

    # ---- 10 核心链路二 ----
    s = blank(prs)
    add_header(s, "五、核心链路二 · 评论 / 表态 / 合集", 10)
    bullets(s, [
        "点赞点踩互斥、收藏幂等（数据库复合主键）",
        "评论两级楼中楼 + 锚点定位 + 图片九宫格",
        "合集 CRUD + 条目拖拽排序 + 笔记合集上下文",
    ], Inches(0.9), Inches(1.6), Inches(4.7), Inches(5), size=15)
    figure(s, FIG / "ch6-fig10-seq-collection.png", Inches(5.7), Inches(1.7), Inches(7.1),
           Inches(4.9), caption="图 6-10  合集侧栏顺序图")

    # ---- 11 详细设计·程序类图 ----
    s = blank(prs)
    add_header(s, "六、详细设计 · 分层与代码对应", 11)
    bullets(s, [
        "routes（drafts / notes / interactions / collections / reports / blocks / auth）→ services → models",
        "文档与真实代码逐项对齐，不写不存在的模块",
    ], Inches(0.9), Inches(1.55), Inches(11.6), Inches(0.85), size=14)
    figure(s, FIG / "ch6-fig1-class-program.png", Inches(0.9), Inches(2.35), Inches(11.5),
           Inches(4.4), caption="图 6-1  程序结构类图")

    # ---- 12 关键技术点 ----
    s = blank(prs)
    add_header(s, "六、关键技术点", 12)
    bullets(s, [
        "滑动窗口限流：按用户计数，互动 30 次 / 10 秒、评论 10 次 / 60 秒",
        "AI 永不硬失败：dry-run / 无 key / 超时 → 自动降级兜底，退避重试",
        "幂等与互斥：复合主键保证幂等，点赞点踩互斥",
        "前端乐观更新：TanStack Query 快照 + 失败回滚",
        "keyset 游标分页：跨页稳定、避免重复 / 漏项",
    ], Inches(0.9), Inches(1.7), Inches(11.6), Inches(5), size=16, gap=13)

    # ---- 13 治理子系统 ----
    s = blank(prs)
    add_header(s, "七、治理子系统（本阶段新增）", 13)
    bullets(s, [
        "举报：建单 / 同人同标的去重 / 内容快照 → 管理员工单队列",
        "AI 后台预审：BackgroundTasks 异步 + 退避重试，高置信自动 ai_flagged",
        "人工裁决：下架 / 删除 / 驳回，保留审计快照",
        "单向拉黑：织入信息流与评论查询过滤，对方无感知",
    ], Inches(0.9), Inches(1.6), Inches(5.1), Inches(5), size=15)
    figure(s, FIG / "ch6-fig13-state-report.png", Inches(6.1), Inches(1.8), Inches(6.7),
           Inches(4.7), caption="图 6-13  举报工单状态图")

    # ---- 14 界面展示 ----
    s = blank(prs)
    add_header(s, "八、界面展示", 14)
    figure(s, FIG / "ch6-fig2-ui-write.png", Inches(0.7), Inches(1.7), Inches(6.0), Inches(4.6),
           caption="写作页：Markdown + AI 差异对照")
    figure(s, FIG / "ch6-fig4-ui-detail.png", Inches(6.9), Inches(1.7), Inches(6.0), Inches(4.6),
           caption="详情页：互动 + 楼中楼评论")

    # ---- 15 测试与验证 ----
    s = blank(prs)
    add_header(s, "九、测试与验证", 15)
    bullets(s, [
        "框架：pytest + httpx.AsyncClient · 内存 SQLite · AI dry-run 零联网",
        "结果：212 passed, 23 skipped",
        "新增测试：操作限流 / 治理(9) / AI 审核(6) / 合集排序",
        "修正的既有问题：签到后 refresh ORM 对象；热门排序按当前周计算",
    ], Inches(0.9), Inches(1.7), Inches(11.6), Inches(5), size=16, gap=13)

    # ---- 16 当前进展与后续 ----
    s = blank(prs)
    add_header(s, "十、当前进展与后续计划", 16)
    bullets(s, [
        "已完成：后端全模块（内容 / 互动 / 成长 / 合集 / 治理 / 限流），212 测试通过",
        "已完成：前端详情页互动、签到成长、合集拖拽、治理入口与管理员工单（开发完成）",
        "进行中：前端本地 typecheck + 构建验收接入",
        "后续：配图导出插入 Word、界面截图完善、AI 审核置信阈值调优",
    ], Inches(0.9), Inches(1.7), Inches(11.6), Inches(5), size=16, gap=13)

    # ---- 17 总结 ----
    s = blank(prs)
    add_header(s, "十一、阶段总结", 17)
    bullets(s, [
        "完成需求 → 概要 → 详细设计 → 测试 的完整闭环，文档与代码逐项对齐",
        "22 张规范化配图：DFD 圆形加工 / 三格类图 / 判断菱形 / 顺序图生命线",
        "后端核心功能全量测试通过（212）；治理子系统后端落地",
        "边界如实：已实现与进行中项在文档与汇报中明确区分",
    ], Inches(0.9), Inches(1.7), Inches(11.6), Inches(5), size=16, gap=13)

    # ---- 18 致谢 ----
    s = blank(prs)
    _rect(s, 0, 0, SW, SH, PRIMARY)
    _, tf = _box(s, Inches(1.0), Inches(2.9), Inches(11.3), Inches(2))
    _para(tf, "谢谢！恳请老师批评指正", 40, WHITE, bold=True, align=PP_ALIGN.CENTER,
          first=True, space_after=18)
    _para(tf, "Q & A", 24, RGBColor(0xBF, 0xD7, 0xEA), align=PP_ALIGN.CENTER, space_after=0)

    prs.save(str(OUT))
    print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build()
